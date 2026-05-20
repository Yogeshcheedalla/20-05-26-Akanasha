import csv
import html
import json
import re
import textwrap
import zipfile
from datetime import datetime
from pathlib import Path
from uuid import uuid4

try:
    import fitz
except Exception:  # pragma: no cover - optional runtime dependency
    fitz = None

try:
    from PIL import Image, ImageDraw, ImageFont
except Exception:  # pragma: no cover - optional runtime dependency
    Image = None
    ImageDraw = None
    ImageFont = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional runtime dependency
    Presentation = None


GENERATED_ARTIFACTS_DIR = Path(__file__).resolve().parents[1] / "generated_artifacts"
GENERATED_ARTIFACTS_DIR.mkdir(parents=True, exist_ok=True)

JAVA_CODING_QUESTIONS = [
    ("Check if a Number is Prime", "Use trial division up to sqrt(n) and reject values below 2.", "public static boolean isPrime(int n) {\n    if (n < 2) return false;\n    for (int i = 2; i * i <= n; i++) {\n        if (n % i == 0) return false;\n    }\n    return true;\n}"),
    ("Factorial Using Iteration", "Multiply all integers from 2 to n; use long for moderate inputs.", "public static long factorial(int n) {\n    long ans = 1;\n    for (int i = 2; i <= n; i++) ans *= i;\n    return ans;\n}"),
    ("Reverse a String", "Walk from the last character to the first using StringBuilder.", "public static String reverse(String s) {\n    StringBuilder out = new StringBuilder();\n    for (int i = s.length() - 1; i >= 0; i--) out.append(s.charAt(i));\n    return out.toString();\n}"),
    ("Palindrome String", "Compare characters from both ends and move inward.", "public static boolean isPalindrome(String s) {\n    int l = 0, r = s.length() - 1;\n    while (l < r) if (s.charAt(l++) != s.charAt(r--)) return false;\n    return true;\n}"),
    ("Fibonacci Series", "Generate terms using two previous values.", "public static int fib(int n) {\n    if (n <= 1) return n;\n    int a = 0, b = 1;\n    for (int i = 2; i <= n; i++) { int c = a + b; a = b; b = c; }\n    return b;\n}"),
    ("Find Maximum in Array", "Track the maximum while scanning once.", "public static int max(int[] arr) {\n    int best = arr[0];\n    for (int x : arr) if (x > best) best = x;\n    return best;\n}"),
    ("Linear Search", "Check every element until target is found.", "public static int linearSearch(int[] arr, int target) {\n    for (int i = 0; i < arr.length; i++) if (arr[i] == target) return i;\n    return -1;\n}"),
    ("Binary Search", "Repeatedly halve a sorted search range.", "public static int binarySearch(int[] arr, int target) {\n    int l = 0, r = arr.length - 1;\n    while (l <= r) {\n        int m = l + (r - l) / 2;\n        if (arr[m] == target) return m;\n        if (arr[m] < target) l = m + 1; else r = m - 1;\n    }\n    return -1;\n}"),
    ("Bubble Sort", "Swap adjacent out-of-order values until sorted.", "public static void bubbleSort(int[] arr) {\n    for (int i = 0; i < arr.length - 1; i++)\n        for (int j = 0; j < arr.length - i - 1; j++)\n            if (arr[j] > arr[j + 1]) { int t = arr[j]; arr[j] = arr[j + 1]; arr[j + 1] = t; }\n}"),
    ("Selection Sort", "Select the minimum remaining element for each position.", "public static void selectionSort(int[] arr) {\n    for (int i = 0; i < arr.length; i++) {\n        int min = i;\n        for (int j = i + 1; j < arr.length; j++) if (arr[j] < arr[min]) min = j;\n        int t = arr[i]; arr[i] = arr[min]; arr[min] = t;\n    }\n}"),
    ("Count Vowels", "Normalize case and count vowel characters.", "public static int countVowels(String s) {\n    int count = 0;\n    for (char c : s.toLowerCase().toCharArray()) if (\"aeiou\".indexOf(c) >= 0) count++;\n    return count;\n}"),
    ("Remove Duplicates from Array", "Use LinkedHashSet to preserve insertion order.", "public static int[] removeDuplicates(int[] arr) {\n    java.util.LinkedHashSet<Integer> set = new java.util.LinkedHashSet<>();\n    for (int x : arr) set.add(x);\n    return set.stream().mapToInt(Integer::intValue).toArray();\n}"),
    ("Frequency of Characters", "Use a map from character to count.", "public static java.util.Map<Character, Integer> frequency(String s) {\n    java.util.Map<Character, Integer> map = new java.util.HashMap<>();\n    for (char c : s.toCharArray()) map.put(c, map.getOrDefault(c, 0) + 1);\n    return map;\n}"),
    ("Anagram Check", "Sort both strings or count characters.", "public static boolean isAnagram(String a, String b) {\n    char[] x = a.toCharArray(), y = b.toCharArray();\n    java.util.Arrays.sort(x); java.util.Arrays.sort(y);\n    return java.util.Arrays.equals(x, y);\n}"),
    ("Second Largest Element", "Track largest and second largest in one pass.", "public static int secondLargest(int[] arr) {\n    int first = Integer.MIN_VALUE, second = Integer.MIN_VALUE;\n    for (int x : arr) { if (x > first) { second = first; first = x; } else if (x > second && x != first) second = x; }\n    return second;\n}"),
    ("Armstrong Number", "Sum digits raised to number of digits.", "public static boolean isArmstrong(int n) {\n    int temp = n, digits = String.valueOf(n).length(), sum = 0;\n    while (temp > 0) { int d = temp % 10; sum += Math.pow(d, digits); temp /= 10; }\n    return sum == n;\n}"),
    ("GCD Using Euclid", "Repeatedly replace larger value with remainder.", "public static int gcd(int a, int b) {\n    while (b != 0) { int t = b; b = a % b; a = t; }\n    return Math.abs(a);\n}"),
    ("LCM of Two Numbers", "Use a * b / gcd(a,b) carefully.", "public static int lcm(int a, int b) {\n    return Math.abs(a / gcd(a, b) * b);\n}"),
    ("Swap Without Third Variable", "Use arithmetic swap when overflow is not a concern.", "public static int[] swap(int a, int b) {\n    a = a + b; b = a - b; a = a - b;\n    return new int[]{a, b};\n}"),
    ("Merge Two Sorted Arrays", "Use two pointers and append the smaller current value.", "public static int[] merge(int[] a, int[] b) {\n    int[] out = new int[a.length + b.length]; int i = 0, j = 0, k = 0;\n    while (i < a.length && j < b.length) out[k++] = a[i] <= b[j] ? a[i++] : b[j++];\n    while (i < a.length) out[k++] = a[i++];\n    while (j < b.length) out[k++] = b[j++];\n    return out;\n}"),
    ("Valid Parentheses", "Use a stack to match closing brackets.", "public static boolean valid(String s) {\n    java.util.Stack<Character> st = new java.util.Stack<>();\n    for (char c : s.toCharArray()) {\n        if (c == '(') st.push(c);\n        else if (c == ')' && (st.empty() || st.pop() != '(')) return false;\n    }\n    return st.empty();\n}"),
    ("Reverse Linked List", "Iteratively reverse next pointers.", "static Node reverse(Node head) {\n    Node prev = null, cur = head;\n    while (cur != null) { Node next = cur.next; cur.next = prev; prev = cur; cur = next; }\n    return prev;\n}"),
    ("Detect Cycle in Linked List", "Use slow and fast pointers.", "static boolean hasCycle(Node head) {\n    Node slow = head, fast = head;\n    while (fast != null && fast.next != null) { slow = slow.next; fast = fast.next.next; if (slow == fast) return true; }\n    return false;\n}"),
    ("Queue Using Two Stacks", "Push into input stack and pop from output stack.", "class MyQueue {\n    java.util.Stack<Integer> in = new java.util.Stack<>(), out = new java.util.Stack<>();\n    void push(int x) { in.push(x); }\n    int pop() { if (out.empty()) while (!in.empty()) out.push(in.pop()); return out.pop(); }\n}"),
    ("Stack Using Queue", "Rotate queue after each push.", "class MyStack {\n    java.util.Queue<Integer> q = new java.util.LinkedList<>();\n    void push(int x) { q.add(x); for (int i = 0; i < q.size() - 1; i++) q.add(q.remove()); }\n    int pop() { return q.remove(); }\n}"),
    ("Find Missing Number", "Use n*(n+1)/2 minus actual sum.", "public static int missing(int[] arr, int n) {\n    int sum = n * (n + 1) / 2;\n    for (int x : arr) sum -= x;\n    return sum;\n}"),
    ("Kadane Maximum Subarray", "Track best subarray ending at current index.", "public static int maxSubArray(int[] arr) {\n    int cur = arr[0], best = arr[0];\n    for (int i = 1; i < arr.length; i++) { cur = Math.max(arr[i], cur + arr[i]); best = Math.max(best, cur); }\n    return best;\n}"),
    ("Two Sum", "Store complement lookup in a hash map.", "public static int[] twoSum(int[] nums, int target) {\n    java.util.Map<Integer, Integer> map = new java.util.HashMap<>();\n    for (int i = 0; i < nums.length; i++) { int need = target - nums[i]; if (map.containsKey(need)) return new int[]{map.get(need), i}; map.put(nums[i], i); }\n    return new int[]{-1, -1};\n}"),
    ("Matrix Diagonal Sum", "Add primary and secondary diagonal values once.", "public static int diagonalSum(int[][] matrix) {\n    int n = matrix.length, sum = 0;\n    for (int i = 0; i < n; i++) { sum += matrix[i][i]; if (i != n - 1 - i) sum += matrix[i][n - 1 - i]; }\n    return sum;\n}"),
    ("String Compression", "Count consecutive repeated characters.", "public static String compress(String s) {\n    StringBuilder out = new StringBuilder();\n    for (int i = 0; i < s.length();) { int j = i; while (j < s.length() && s.charAt(j) == s.charAt(i)) j++; out.append(s.charAt(i)).append(j - i); i = j; }\n    return out.toString();\n}"),
]


def requested_artifact_formats(prompt: str) -> list[str]:
    lowered = prompt.lower()
    if re.search(r"\b(all\s+(?:file\s+)?formats?|all\s+types\s+of\s+(?:files?|documents?|formats?))\b", lowered):
        return ["pdf", "docx", "pptx", "xlsx", "csv", "json", "png", "jpg", "md", "zip"]
    formats: list[str] = []
    checks = [
        ("xlsx", r"\b(excel|xlsx|spreadsheets?|workbooks?)\b"),
        ("pdf", r"\b(pdfs?|reports?|invoices?|receipts?|certificates?|resume|resumes|notes?|formula sheets?|study plans?)\b"),
        ("docx", r"\b(word|docx|documents?)\b"),
        ("pptx", r"\b(powerpoints?|ppt|pptx|presentations?|slides?)\b"),
        ("csv", r"\b(csvs?|csv files?)\b"),
        ("json", r"\b(json)\b"),
        ("png", r"\b(pngs?|images?|diagrams?|photos?|pictures?|charts?)\b"),
        ("jpg", r"\b(jpgs?|jpegs?|jpcs?)\b"),
        ("md", r"\b(markdown|md)\b"),
        ("zip", r"\b(zips?|archives?)\b"),
    ]
    for name, pattern in checks:
        if re.search(pattern, lowered):
            formats.append(name)
    return formats


def sanitize_model_artifact_placeholders(response_text: str) -> str:
    """Remove model-invented download links before appending real generated files.

    The LLM sometimes emits links such as sandbox:/file.pdf even though only the
    backend artifact engine can create real downloadable files in this app.
    """
    cleaned = re.sub(r"\[[^\]]+\]\(\s*sandbox:[^)]+\)", "", response_text, flags=re.IGNORECASE)
    cleaned = re.sub(r"(?im)^\s*(?:download|open)\s+[^:\n]*:\s*sandbox:[^\n]*$", "", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def _safe_slug(prompt: str) -> str:
    words = re.findall(r"[A-Za-z0-9]+", prompt.lower())[:8]
    return "-".join(words) or "akansha-output"


def _artifact_path(prompt: str, extension: str) -> Path:
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return GENERATED_ARTIFACTS_DIR / f"{_safe_slug(prompt)}-{stamp}-{uuid4().hex[:8]}.{extension}"


def _strip_markdown(value: str) -> str:
    value = re.sub(r"```[A-Za-z0-9_+\-.#]*\s*\n([\s\S]*?)```", r"\1", value)
    value = re.sub(r"```([\s\S]*?)```", r"\1", value)
    value = re.sub(r"`([^`]+)`", r"\1", value)
    value = re.sub(r"\*\*([^*]+)\*\*", r"\1", value)
    return value.strip()


def _requested_count(prompt: str, nouns: tuple[str, ...]) -> int | None:
    noun_pattern = "|".join(re.escape(noun) for noun in nouns)
    patterns = [
        rf"\b(?:at least|min(?:imum)?|minimum of)?\s*(\d{{1,3}})\s+(?:{noun_pattern})\b",
        rf"\b(?:{noun_pattern})\s*(?:of|:)?\s*(\d{{1,3}})\b",
    ]
    for pattern in patterns:
        match = re.search(pattern, prompt, flags=re.IGNORECASE)
        if match:
            return max(1, min(int(match.group(1)), 120))
    return None


def _is_java_coding_pack_request(prompt: str) -> bool:
    lowered = prompt.lower()
    return "java" in lowered and bool(re.search(r"\b(code|coding|program|programs|questions?|placements?|tcs)\b", lowered))


def _build_java_coding_pack(prompt: str, response_text: str) -> str:
    question_count = _requested_count(prompt, ("questions?", "coding questions?", "programs?", "problems?")) or 30
    question_count = max(10, min(question_count, len(JAVA_CODING_QUESTIONS)))
    title = "Java Coding Questions for TCS Placements"
    lines = [
        f"# {title}",
        "",
        "Purpose: a practice-ready Java coding pack with complete code, comments, approach, complexity, and sample tests.",
        "Language: Java.",
        f"Question count: {question_count}.",
        "",
        "## How to Use This Pack",
        "- Read the problem statement first.",
        "- Understand the approach before memorizing code.",
        "- Dry-run the sample input/output.",
        "- Recode the solution without looking.",
        "- Revise time and space complexity before interviews.",
        "",
        "## Topic Coverage",
        "- Basic: numbers, strings, arrays, search, sorting.",
        "- Medium: hashing, two pointers, stack, queue, linked list.",
        "- Advanced: Kadane, matrix, cycle detection, data-structure design.",
        "",
    ]
    for index, (name, approach, code) in enumerate(JAVA_CODING_QUESTIONS[:question_count], start=1):
        level = "Basic" if index <= 10 else "Medium" if index <= 22 else "Advanced"
        lines.extend(
            [
                f"## Question {index}: {name}",
                f"Level: {level}",
                "",
                f"Problem: Write a Java program to solve: {name.lower()}.",
                f"Approach: {approach}",
                "",
                "Complete Java code with comments:",
                "```java",
                "import java.util.*;",
                "",
                "public class Solution {",
                *[f"    {line}" if line else "" for line in code.splitlines()],
                "",
                "    public static void main(String[] args) {",
                f"        // Practice entry point for: {name}",
                "        // Add scanner input here according to the platform format.",
                "    }",
                "}",
                "```",
                "",
                "Sample test idea: run with a normal case, an edge case, and a repeated/large-value case.",
                "Time Complexity: usually O(n), O(log n), or O(n log n) depending on the algorithm above.",
                "Space Complexity: O(1) unless an extra map, set, stack, queue, or output array is used.",
                "Interview Tip: explain why the algorithm works before writing the code.",
                "",
            ]
        )
    lines.extend(
        [
            "## Final Revision Checklist",
            "- Can you explain every loop condition?",
            "- Can you state time and space complexity?",
            "- Did you handle null, empty, one-element, and duplicate cases?",
            "- Can you convert the logic to scanner-based input if TCS asks for stdin/stdout?",
        ]
    )
    if response_text and "sandbox:" not in response_text.lower():
        lines.extend(["", "## Assistant Draft Notes", _strip_markdown(response_text)[:1200]])
    return "\n".join(lines)


def _build_invoice_content(prompt: str, response_text: str) -> str:
    if not re.search(r"\binvoice\b", prompt, flags=re.IGNORECASE):
        return response_text
    amount = re.search(r"(?:₹|rs\.?|inr)\s*([0-9,]+)", prompt, flags=re.IGNORECASE)
    client = re.search(r"\bclient(?: name)?\s+([A-Za-z][A-Za-z\s]{1,40}?)(?:,| amount| for|$)", prompt, flags=re.IGNORECASE)
    service = re.search(r"\bfor\s+([^,.]+)", prompt, flags=re.IGNORECASE)
    return "\n".join(
        [
            "# Invoice",
            "",
            f"Client: {client.group(1).strip() if client else 'Client'}",
            f"Service: {service.group(1).strip() if service else 'Professional service'}",
            f"Amount: INR {amount.group(1) if amount else '0'}",
            f"Generated: {datetime.now().strftime('%d %b %Y, %I:%M %p')}",
            "",
            "| Item | Description | Amount |",
            "|---|---|---:|",
            f"| 1 | {service.group(1).strip() if service else 'Professional service'} | {amount.group(1) if amount else '0'} |",
            "",
            "Payment Status: Pending",
            "Notes: This invoice was generated by Akansha from the user's prompt.",
            "",
            _strip_markdown(response_text)[:1000],
        ]
    )


def _prepare_artifact_content(prompt: str, response_text: str) -> str:
    cleaned = sanitize_model_artifact_placeholders(response_text)
    if _is_java_coding_pack_request(prompt):
        return _build_java_coding_pack(prompt, cleaned)
    if re.search(r"\binvoice\b", prompt, flags=re.IGNORECASE):
        return _build_invoice_content(prompt, cleaned)
    return cleaned


def _extract_markdown_table(text: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not (stripped.startswith("|") and stripped.endswith("|")):
            continue
        cells = [cell.strip() for cell in stripped.strip("|").split("|")]
        if cells and all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells):
            continue
        rows.append(cells)
    return rows


def _rows_from_response(response_text: str) -> list[list[str]]:
    table_rows = _extract_markdown_table(response_text)
    if table_rows:
        return table_rows

    rows = [["Section", "Content"]]
    for index, line in enumerate(_strip_markdown(response_text).splitlines(), start=1):
        cleaned = line.strip(" -\t")
        if cleaned:
            rows.append([str(index), cleaned])
    return rows if len(rows) > 1 else [["Content"], [_strip_markdown(response_text)]]


def _write_csv(path: Path, response_text: str) -> None:
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.writer(handle)
        writer.writerows(_rows_from_response(response_text))


def _xml_escape(value: str) -> str:
    return html.escape(value, quote=True)


def _write_xlsx(path: Path, response_text: str) -> None:
    rows = _rows_from_response(response_text)
    sheet_rows = []
    for row_index, row in enumerate(rows, start=1):
        cells = []
        for column_index, value in enumerate(row, start=1):
            column_name = ""
            current = column_index
            while current:
                current, remainder = divmod(current - 1, 26)
                column_name = chr(65 + remainder) + column_name
            cell_ref = f"{column_name}{row_index}"
            cells.append(
                f'<c r="{cell_ref}" t="inlineStr"><is><t>{_xml_escape(str(value))}</t></is></c>'
            )
        sheet_rows.append(f'<row r="{row_index}">{"".join(cells)}</row>')

    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
<Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/workbook.xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
<sheets><sheet name="Akansha Output" sheetId="1" r:id="rId1"/></sheets>
</workbook>""",
        )
        archive.writestr(
            "xl/_rels/workbook.xml.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "xl/worksheets/sheet1.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
<sheetData>{"".join(sheet_rows)}</sheetData>
</worksheet>""",
        )


def _write_pdf(path: Path, response_text: str, prompt: str = "") -> None:
    if fitz is None:
        path.write_text(_strip_markdown(response_text), encoding="utf-8")
        return
    min_pages = _requested_count(prompt, ("pages?",)) or 1
    document = fitz.open()
    page = document.new_page(width=595, height=842)
    page_count = 1
    y = 54
    page.insert_text((48, y), "Akansha Report", fontsize=18, fontname="helv", color=(0.18, 0.12, 0.42))
    y += 34
    for paragraph in _strip_markdown(response_text).splitlines():
        for wrapped in textwrap.wrap(paragraph, width=88) or [""]:
            if y > 790:
                page = document.new_page(width=595, height=842)
                page_count += 1
                y = 54
            page.insert_text((48, y), wrapped, fontsize=10, fontname="helv")
            y += 16
    while page_count < min_pages:
        page = document.new_page(width=595, height=842)
        page_count += 1
        page.insert_text((48, 54), f"Practice Page {page_count}", fontsize=16, fontname="helv", color=(0.18, 0.12, 0.42))
        y = 90
        for line in [
            "Use this page for dry runs, handwritten logic, complexity notes, and revision.",
            "Problem:",
            "Approach:",
            "Code notes:",
            "Edge cases:",
            "Time complexity:",
            "Space complexity:",
        ]:
            page.insert_text((48, y), line, fontsize=11, fontname="helv")
            y += 28
    document.save(path)
    document.close()


def _write_docx(path: Path, response_text: str) -> None:
    paragraphs = "".join(
        f"<w:p><w:r><w:t>{_xml_escape(line)}</w:t></w:r></w:p>"
        for line in _strip_markdown(response_text).splitlines()
        if line.strip()
    )
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>""",
        )
        archive.writestr(
            "word/document.xml",
            f"""<?xml version="1.0" encoding="UTF-8"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body>{paragraphs}<w:sectPr/></w:body>
</w:document>""",
        )


def _write_pptx(path: Path, response_text: str) -> None:
    if Presentation is None:
        path.write_text(_strip_markdown(response_text), encoding="utf-8")
        return
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Akansha Report"
    title_slide.placeholders[1].text = datetime.now().strftime("%d %b %Y, %I:%M %p")
    lines = [line.strip(" -") for line in _strip_markdown(response_text).splitlines() if line.strip()]
    for chunk_start in range(0, len(lines), 7):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Key points"
        body = slide.placeholders[1].text_frame
        body.clear()
        for line in lines[chunk_start : chunk_start + 7]:
            paragraph = body.add_paragraph()
            paragraph.text = line[:180]
            paragraph.level = 0
    presentation.save(path)


def _write_image(path: Path, response_text: str, image_format: str) -> None:
    if Image is None or ImageDraw is None:
        path.write_text(_strip_markdown(response_text), encoding="utf-8")
        return
    width, height = 1400, 900
    image = Image.new("RGB", (width, height), "#080914")
    draw = ImageDraw.Draw(image)
    font_title = ImageFont.load_default()
    font_body = ImageFont.load_default()
    draw.rectangle((0, 0, width, 110), fill="#17142d")
    draw.text((52, 40), "Akansha Output", fill="#ffffff", font=font_title)
    y = 145
    for line in _strip_markdown(response_text).splitlines():
        for wrapped in textwrap.wrap(line, width=110) or [""]:
            if y > height - 60:
                break
            draw.text((52, y), wrapped, fill="#f4f0ff", font=font_body)
            y += 24
    image.save(path, "JPEG" if image_format == "jpg" else "PNG", quality=94)


def create_requested_artifacts(prompt: str, response_text: str) -> list[dict[str, str]]:
    formats = requested_artifact_formats(prompt)
    if not formats:
        return []
    artifact_content = _prepare_artifact_content(prompt, response_text)
    if "zip" in formats and len(formats) == 1:
        formats = ["md", "json", "csv", "zip"]

    created: list[dict[str, str]] = []
    generated_paths: list[Path] = []
    for fmt in formats:
        if fmt == "zip":
            continue
        path = _artifact_path(prompt, fmt)
        if fmt == "xlsx":
            _write_xlsx(path, artifact_content)
        elif fmt == "pdf":
            _write_pdf(path, artifact_content, prompt)
        elif fmt == "docx":
            _write_docx(path, artifact_content)
        elif fmt == "pptx":
            _write_pptx(path, artifact_content)
        elif fmt == "csv":
            _write_csv(path, artifact_content)
        elif fmt == "json":
            path.write_text(
                json.dumps({"prompt": prompt, "content": artifact_content, "created_at": datetime.now().isoformat()}, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        elif fmt == "md":
            path.write_text(artifact_content, encoding="utf-8")
        elif fmt in {"png", "jpg"}:
            _write_image(path, artifact_content, fmt)
        else:
            continue
        generated_paths.append(path)
        created.append({"format": fmt, "name": path.name, "url": f"/generated/{path.name}"})

    if "zip" in formats and generated_paths:
        zip_path = _artifact_path(prompt, "zip")
        with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for generated_path in generated_paths:
                archive.write(generated_path, generated_path.name)
        created.append({"format": "zip", "name": zip_path.name, "url": f"/generated/{zip_path.name}"})

    return created


def artifact_markdown(artifacts: list[dict[str, str]]) -> str:
    if not artifacts:
        return ""
    rows = ["| Format | Download |", "|---|---|"]
    for artifact in artifacts:
        rows.append(f"| {artifact['format'].upper()} | [{artifact['name']}]({artifact['url']}) |")
    return "\n\nGenerated files:\n" + "\n".join(rows)
